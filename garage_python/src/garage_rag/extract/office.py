"""Office document extraction: .docx, .pptx, .xlsx.

Legacy binary formats (.doc, .ppt, .xls) are not handled -- none of these
libraries read them. The dispatcher reports them as unsupported rather than
silently producing empty documents; converting them needs LibreOffice.
"""

from __future__ import annotations

import logging
from pathlib import Path

from .base import (
    ContentKind,
    ExtractionError,
    ExtractResult,
    clean_author_hints,
    normalize_text,
)

log = logging.getLogger(__name__)

VERSION = "1"

# Spreadsheets can be enormous; a runaway sheet should not blow up memory.
MAX_SHEET_ROWS = 5000
MAX_CELL_CHARS = 500


def _core_properties(props) -> tuple[dict, list[str], str | None]:  # noqa: ANN001
    """Shared docx/pptx core-properties handling."""
    meta: dict = {}
    hints: list[str] = []
    title: str | None = None

    try:
        author = (props.author or "").strip()
        if author:
            meta["office_author"] = author[:500]
            hints.append(author)
        last_by = (props.last_modified_by or "").strip()
        if last_by and last_by != author:
            meta["office_last_modified_by"] = last_by[:500]
            hints.append(last_by)
        title = (props.title or "").strip() or None
        if title:
            meta["office_title"] = title[:500]
    except Exception as exc:  # noqa: BLE001
        log.debug("core properties unreadable: %s", exc)

    return meta, hints, title


def extract_docx(path: Path) -> ExtractResult:
    import docx

    try:
        document = docx.Document(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"unreadable .docx {path}: {exc}") from exc

    meta, hints, title = _core_properties(document.core_properties)

    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        # Preserve heading structure as Markdown so the markdown-aware chunker
        # can use it for heading paths.
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit())
            parts.append(f"{'#' * (int(level) if level else 1)} {text}")
        else:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    text = normalize_text("\n\n".join(parts))
    if not text:
        raise ExtractionError(f"no text extracted from {path}")

    return ExtractResult(
        text=text,
        kind=ContentKind.MARKDOWN,
        extractor="python-docx",
        extractor_version=VERSION,
        title=title or path.stem,
        meta=meta,
        author_hints=clean_author_hints(hints),
    )


def extract_pptx(path: Path) -> ExtractResult:
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"unreadable .pptx {path}: {exc}") from exc

    meta, hints, title = _core_properties(presentation.core_properties)

    parts: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        slide_parts.append(line)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        slide_parts.append(" | ".join(cells))

        # Speaker notes often hold the actual argument, not just the bullets.
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_parts.append(f"Notes: {notes}")
        except Exception:  # noqa: BLE001
            pass

        if slide_parts:
            # Slide headings give the chunker natural boundaries.
            parts.append(f"## Slide {number}\n" + "\n".join(slide_parts))

    text = normalize_text("\n\n".join(parts))
    if not text:
        raise ExtractionError(f"no text extracted from {path}")

    meta["slide_count"] = len(presentation.slides)
    return ExtractResult(
        text=text,
        kind=ContentKind.MARKDOWN,
        extractor="python-pptx",
        extractor_version=VERSION,
        title=title or path.stem,
        meta=meta,
        author_hints=clean_author_hints(hints),
    )


def extract_xlsx(path: Path) -> ExtractResult:
    import openpyxl

    try:
        # read_only + data_only: stream rows, and take cached formula results
        # rather than the formula text, which is what a reader would see.
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"unreadable .xlsx {path}: {exc}") from exc

    meta: dict = {}
    hints: list[str] = []
    try:
        creator = (workbook.properties.creator or "").strip()
        if creator:
            meta["office_author"] = creator[:500]
            hints.append(creator)
    except Exception:  # noqa: BLE001
        pass

    parts: list[str] = []
    truncated: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= MAX_SHEET_ROWS:
                    truncated.append(sheet.title)
                    break
                cells = [str(value)[:MAX_CELL_CHARS] for value in row if value is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"## {sheet.title}\n" + "\n".join(rows))
    finally:
        workbook.close()

    text = normalize_text("\n\n".join(parts))
    if not text:
        raise ExtractionError(f"no text extracted from {path}")

    if truncated:
        meta["truncated_sheets"] = truncated
    meta["sheet_count"] = len(parts)

    return ExtractResult(
        text=text,
        kind=ContentKind.TABULAR,
        extractor="openpyxl",
        extractor_version=VERSION,
        title=path.stem,
        meta=meta,
        author_hints=clean_author_hints(hints),
    )
